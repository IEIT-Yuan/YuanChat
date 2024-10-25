"""
A set of functions to create a PowerPoint slide deck.
"""
import time
import random
import re
import sys
import io
from typing import List, Tuple, Optional
import pptx
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches,Pt,Cm
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR,PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
sys.path.append('..')
sys.path.append('../..')
import os
from pkg.ppt_gen import template_helper as template
from pkg.llm import generate_ppt_stream
from pkg.ppt_gen import picture_helper as ims
from pkg.ppt_gen import config

from concurrent.futures import ThreadPoolExecutor



current_work_path = os.path.join(os.getcwd(), "icons", "template_background","青春.pptx")

load_dotenv()

# English Metric Unit (used by PowerPoint) to inches
EMU_TO_INCH_SCALING_FACTOR = 1.0 / 914400
INCHES_3 = Inches(3)
INCHES_2 = Inches(2)
INCHES_1_5 = Inches(1.5)
INCHES_1 = Inches(1)
INCHES_0_8 = Inches(0.8)
INCHES_0_9 = Inches(0.9)
INCHES_0_5 = Inches(0.5)
INCHES_0_4 = Inches(0.4)
INCHES_0_3 = Inches(0.3)
INCHES_0_2 = Inches(0.2)
ICON_SIZE = INCHES_0_8
ICON_BG_SIZE = INCHES_1

IMAGE_DISPLAY_PROBABILITY = 1 / 3.0
FOREGROUND_IMAGE_PROBABILITY = 0.8

SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
ICONS_REGEX = re.compile(r"\[\[(.*?)\]\]\s*(.*)")
icon_path=os.path.join(os.getcwd(), "icons", "png128",'value.png')
def get_flat_list_of_contents(items: list, level: int) -> List[Tuple]:
    """
    Flatten a (hierarchical) list of bullet points to a single list containing each item and
    its level.

    :param items: A bullet point (string or list).
    :param level: The current level of hierarchy.
    :return: A list of (bullet item text, hierarchical level) tuples.
    """

    flat_list = []

    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list = flat_list + get_flat_list_of_contents(item, level + 1)

    return flat_list
def _add_text_at_bottom(
        slide: pptx.slide.Slide,
        slide_width_inch: float,
        slide_height_inch: float,
        text: str,
        hyperlink: Optional[str] = None,
        target_height: Optional[float] = 0.5
):
    """
    Add arbitrary text to a textbox positioned near the lower left side of a slide.

    :param slide: The slide.
    :param slide_width_inch: The width of the slide.
    :param slide_height_inch: The height of the slide.
    :param target_height: the target height of the box in inches (optional).
    :param text: The text to be added
    :param hyperlink: The hyperlink to be added to the text (optional).
    """

    footer = slide.shapes.add_textbox(
        left=INCHES_1,
        top=pptx.util.Inches(slide_height_inch - target_height),
        width=pptx.util.Inches(slide_width_inch),
        height=pptx.util.Inches(target_height)
    )
    footer.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    paragraph = footer.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = pptx.util.Pt(10)
    run.font.underline = False

    if hyperlink:
        run.hyperlink.address = hyperlink
def _get_slide_width_height_inches(presentation: pptx.Presentation) -> Tuple[float, float]:
    """
    Get the dimensions of a slide in inches.

    :param presentation: The presentation object.
    :return: The width and the height.
    """

    slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_width
    slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_height
    # logger.debug('Slide width: %f, height: %f', slide_width_inch, slide_height_inch)

    return slide_width_inch, slide_height_inch

#排版左图右文字的幻灯片
def add_signal_content_slide(ppt:pptx.Presentation,data:dict,title:str,img,img_Dict:dict,key_message:str,title_color:RGBColor):
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    original_left = 0
    original_top = 0
    left_width = 0
    left_height = 0
    #插入图片到幻灯片左边三分之一处
    if img:
        picture = slide.shapes.add_picture(img,left=Cm(0),top=Cm(0),width=ppt.slide_width/3,height=ppt.slide_height)
    #添加标题文本框
    title_box = slide.shapes.add_textbox(left=ppt.slide_width/3+ppt.slide_width/20,top=Cm(0),width=ppt.slide_width*2/3,height=ppt.slide_height/5)
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_box.text_frame.word_wrap = True
    title_box.text_frame.paragraphs[0].text = title
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.color.rgb = title_color
    #添加key_message
    key_message_box = slide.shapes.add_textbox(left=ppt.slide_width/3+ppt.slide_width/20,top=ppt.slide_height*3/25+ppt.slide_height/50,width=ppt.slide_width*2/3,height=ppt.slide_height/25)
    key_message_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    key_message_box.text_frame.word_wrap = True
    key_message_box.text_frame.paragraphs[0].text = key_message
    key_message_box.text_frame.paragraphs[0].font.size = Pt(10) 
    key_message_box.text_frame.paragraphs[0].font.color.rgb = title_color
    #添加装饰
    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER,left=ppt.slide_width/3,top=ppt.slide_height/6+ppt.slide_height/20,width=ppt.slide_width*2/3,height=Cm(0.17))
    shape.fill.solid()
    shape.fill.fore_color.rgb = title_color
    shape.shadow.inherit=True
    shape.line.fill.background()
    #添加内容文本框
    original_left = ppt.slide_width/3+ppt.slide_width/20
    original_top = ppt.slide_height/5+ppt.slide_height/10
    left_height = ppt.slide_height*5/6
    left_width = ppt.slide_width*2/3-ppt.slide_width/10
    
    per_height = left_height/(len(data)+1)
    perper_height = per_height/(len(data)+1)
    for k,v in data.items():
        content_title_box = slide.shapes.add_textbox(left = original_left,top = original_top,width = left_width,height = perper_height*2)
        content_box = slide.shapes.add_textbox(left = original_left,top = original_top+perper_height*12/10,width=left_width,height =perper_height)
        content_title_box.text_frame.paragraphs[0].text = k
        content_box.text_frame.word_wrap = True
        content_box.text_frame.paragraphs[0].text = v
        content_box.text_frame.paragraphs[0].font.size = Pt(10)
        content_title_box.text_frame.paragraphs[0].font.bold = True
        content_title_box.text_frame.paragraphs[0].font.color.rgb = title_color
        content_title_box.text_frame.paragraphs[0].font.size = Pt(14)
        original_top+=per_height

    #添加背景
    for k,v in img_Dict.items():
        picture = slide.shapes.add_picture(k,v[0],v[1],v[2],v[3])
        slide.shapes._spTree.remove(picture._element)
        slide.shapes._spTree.insert(2, picture._element)
    return True

#排版内容要点为四个，无图片插入的幻灯片
def add_four_content_slide_no_picture(ppt:pptx.Presentation,data:dict,title:str,img_Dict:dict,key_message:str,title_color:RGBColor):
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    #插入标题
    title_box = slide.shapes.add_textbox(left=ppt.slide_width/20,top=ppt.slide_height/20,width=ppt.slide_width*19/20,height=ppt.slide_height/10)
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_box.text_frame.paragraphs[0].text = title
    title_box.text_frame.paragraphs[0].font.bold = True
    #添加key_message
    key_message_box = slide.shapes.add_textbox(left=ppt.slide_width/20,top=ppt.slide_height/20+ppt.slide_height/10,width=ppt.slide_width*19/20,height=ppt.slide_height/20)
    key_message_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    key_message_box.text_frame.word_wrap = True
    key_message_box.text_frame.paragraphs[0].text = key_message
    key_message_box.text_frame.paragraphs[0].font.size = Pt(10) 
    key_message_box.text_frame.paragraphs[0].font.color.rgb = title_color
    #渐变色背景
    title_box.fill.gradient()
    title_box.fill.gradient_angle = 45
    title_box.fill.gradient_stops[0].color.rgb = title_color
    title_box.fill.gradient_stops[1].color.rgb = RGBColor(255,255,255)

    #插入装饰线
    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER,left=ppt.slide_width/20,top=ppt.slide_height/4,width=ppt.slide_width*9/10,height=ppt.slide_height/400)
    shape.fill.solid()
    shape.fill.fore_color.rgb = title_color
    shape.shadow.inherit=True
    shape.line.fill.background()

    #添加内容框
    page = len(data)
    line = ppt.slide_width*9/10
    interval = line/(page+1)
    original_left = ppt.slide_width/20
    original_top = ppt.slide_height*7/20
    need_width = interval*9/10
    need_height = ppt.slide_height/2
    count = 1
    for k,v in data.items():
        content_box = slide.shapes.add_textbox(left = original_left+interval-need_width/2+(count-1)*interval,top = original_top+need_height/3,width = need_width,height = need_height*2/3)
        title_box = slide.shapes.add_textbox(left = original_left+interval-need_width/2+(count-1)*interval,top = original_top,width = need_width,height = need_height/3)
        #内部标题
        title_box.text_frame.word_wrap = True
        title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_box.text_frame.paragraphs[0].text = k
        title_box.text_frame.paragraphs[0].font.size = Pt(14)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        #内部主体内容
        content_box.text_frame.word_wrap = True
        content_box.text_frame.paragraphs[0].line_spacing = 1.5
        content_box.text_frame.paragraphs[0].text = v
        content_box.text_frame.paragraphs[0].font.size = Pt(10)
        #添加背景颜色
        # content_box.fill.solid()
        # content_box.fill.fore_color.rgb = RGBColor(242,242,242)
        #添加点装饰
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,left=original_left+interval*count,top=ppt.slide_height/4-ppt.slide_height/200+ppt.slide_height/800,width=ppt.slide_height/100,height=ppt.slide_height/100)
        shape.fill.solid()
        shape.fill.fore_color.rgb = title_color
        shape.shadow.inherit=True
        shape.line.fill.background()
        #添加序号按钮
        shape = slide.shapes.add_shape(MSO_SHAPE.PLAQUE,left=original_left+interval*count-ppt.slide_width/32,top=original_top-ppt.slide_height/40,width=ppt.slide_width/16,height=ppt.slide_height/20)
        shape.text_frame.paragraphs[0].text = str(count)
        shape.text_frame.paragraphs[0].font.color.rgb = title_color
        shape.fill.background()
        shape.line.color.rgb = title_color
        #下一个框
        count+=1
    #添加背景
    for k,v in img_Dict.items():
        picture = slide.shapes.add_picture(k,v[0],v[1],v[2],v[3])
        slide.shapes._spTree.remove(picture._element)
        slide.shapes._spTree.insert(2, picture._element)
    return True

#排版带单张图片且有四个内容的幻灯片
def add_four_content_slide_with_signal_picture(ppt:pptx.Presentation,data:dict,img_key,icon_path:str,title:str,img_Dict:dict,key_message:str,title_color:RGBColor):
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    #添加标题
    title_box = slide.shapes.add_textbox(left=ppt.slide_width/20,top=ppt.slide_height/20,width=ppt.slide_width*19/20,height=ppt.slide_height/10)
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_box.text_frame.paragraphs[0].text = title
    title_box.text_frame.paragraphs[0].font.bold = True
    #添加key_message
    key_message_box = slide.shapes.add_textbox(left=ppt.slide_width/20,top=ppt.slide_height/20+ppt.slide_height/10,width=ppt.slide_width*19/20,height=ppt.slide_height/20)
    key_message_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    key_message_box.text_frame.word_wrap = True
    key_message_box.text_frame.paragraphs[0].text = key_message
    key_message_box.text_frame.paragraphs[0].font.size = Pt(10) 
    key_message_box.text_frame.paragraphs[0].font.color.rgb = title_color
    #渐变色背景
    title_box.fill.gradient()
    title_box.fill.gradient_angle = 45
    title_box.fill.gradient_stops[0].color.rgb = title_color
    title_box.fill.gradient_stops[1].color.rgb = RGBColor(255,255,255)
    #添加图片
    if img_key:
        picture = slide.shapes.add_picture(img_key,left=ppt.slide_width*11/30,top=ppt.slide_height*3/8-ppt.slide_height/20,width=ppt.slide_width*4/15,height=ppt.slide_height/2)
    #添加文本框
    count =0
    hang = 0
    lie = 0
    for k,v in data.items():
        content_box = slide.shapes.add_textbox(left = ppt.slide_width/30+hang*ppt.slide_width*2/3,top = ppt.slide_height*3/8-ppt.slide_height/20+lie*ppt.slide_height/4+lie*(ppt.slide_height/8),width = ppt.slide_width*4/15,height = ppt.slide_height/4)
        #内部标题

        content_box.text_frame.word_wrap = True
        content_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        content_box.text_frame.paragraphs[0].text = k
        content_box.text_frame.paragraphs[0].font.size = Pt(14)
        content_box.text_frame.paragraphs[0].font.bold = True
        if hang==0:
            content_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        else:
            content_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        content_box.text_frame.paragraphs[0].font.color.rgb = title_color
        #内部主体内容
        para = content_box.text_frame.add_paragraph()
        para.line_spacing = 1.5
        para.text = v
        para.font.size = Pt(10)
        content_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        #添加背景颜色
        # content_box.fill.solid()
        # content_box.fill.fore_color.rgb = RGBColor(242,242,242)
        #对应位置加Icon
        pic = slide.shapes.add_picture(icon_path,left=ppt.slide_width*13/48+hang*ppt.slide_width*19/48,top=ppt.slide_height*5/16-ppt.slide_height/20+lie*ppt.slide_height*3/8,width=ppt.slide_width/16,height=ppt.slide_height/16)
        count+=1
        hang = hang+1
        if count==2:
            hang = 0
            lie = 1
        elif count == 3:
            hang = 1
            lie = 1
    #添加下修饰线
    shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER,left=ppt.slide_width*11/30,top=99*ppt.slide_height/100,width=ppt.slide_width*4/15,height=ppt.slide_height/100)
    shape.line.color.rgb = title_color
    shape.shadow.inherit=True
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = title_color

    #添加背景
    for k,v in img_Dict.items():
        picture = slide.shapes.add_picture(k,v[0],v[1],v[2],v[3])
        slide.shapes._spTree.remove(picture._element)
        slide.shapes._spTree.insert(2, picture._element)

#横向排版带多张图片且有多个内容的幻灯片
def add_content_slide_with_multiple_picture(ppt:pptx.Presentation,data:dict,img_key:list,title:str,img_Dict:dict,key_message:str,title_color:RGBColor):
    page = len(data) #data包含的长度
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    #添加标题
    title_box = slide.shapes.add_textbox(left=ppt.slide_width/20,top=ppt.slide_height/20,width=ppt.slide_width*19/20,height=ppt.slide_height/10)
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_box.text_frame.paragraphs[0].text = title
    title_box.text_frame.paragraphs[0].font.bold = True
    #添加key_message
    key_message_box = slide.shapes.add_textbox(left=ppt.slide_width/20,top=ppt.slide_height/20+ppt.slide_height/10,width=ppt.slide_width*19/20,height=ppt.slide_height/20)
    key_message_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    key_message_box.text_frame.word_wrap = True
    key_message_box.text_frame.paragraphs[0].text = key_message
    key_message_box.text_frame.paragraphs[0].font.size = Pt(10) 
    key_message_box.text_frame.paragraphs[0].font.color.rgb = title_color
    #渐变色背景
    title_box.fill.gradient()
    title_box.fill.gradient_angle = 45
    title_box.fill.gradient_stops[0].color.rgb = title_color
    title_box.fill.gradient_stops[1].color.rgb = RGBColor(255,255,255)
    #插入装饰线
    # shape = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER,left=ppt.slide_width/20,top=ppt.slide_height/4,width=ppt.slide_width*9/10,height=Cm(0.05))
    # shape.line.color.rgb = RGBColor(71,132,203)
    # shape.shadow.inherit=True
    # shape.line.fill.background()
    #添加内容
    #高度固定 只变宽度

    shape_line = ppt.slide_width*9/10
    interval = shape_line/10  #间隔

    per_interval = interval/(page-1)
    count = 0 
    org_height = ppt.slide_height/4
    text_height = ppt.slide_height*2/3
    percontent_width =shape_line*9/10/page
    org_width = ppt.slide_width/20
    content_width = percontent_width
    count_content = 0
    count_interval=0
    for k,v in data.items():
        #添加图片
        if count<len(img_key):
            if img_key[count]:
                picture = slide.shapes.add_picture(img_key[count],left=org_width+count_content*percontent_width+count_interval*per_interval,top=org_height*9/10,width=content_width,height=content_width)
        #添加文本框
        content_box = slide.shapes.add_textbox(left = org_width+count_content*percontent_width+count_interval*per_interval,top = org_height+content_width,width = content_width,height = ppt.slide_height-(org_height+content_width))
        #添加标题
        content_box.text_frame.word_wrap = True
        # content_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        content_box.text_frame.paragraphs[0].text = k
        content_box.text_frame.paragraphs[0].font.size = Pt(14)
        content_box.text_frame.paragraphs[0].font.bold = True#加粗
        content_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER#居中
        content_box.text_frame.paragraphs[0].font.color.rgb = title_color
        para = content_box.text_frame.add_paragraph()
        para.line_spacing = 1.5
        para.text = v
        para.font.size = Pt(10)
        content_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        #添加背景颜色
        # content_box.fill.solid()
        # content_box.fill.fore_color.rgb = RGBColor(242,242,242)
        count_content+=1
        count_interval+=1
        count+=1

    #添加背景
    for k,v in img_Dict.items():
        picture = slide.shapes.add_picture(k,v[0],v[1],v[2],v[3])
        slide.shapes._spTree.remove(picture._element)
        slide.shapes._spTree.insert(2, picture._element)

    return True

#排版单个带图片的幻灯片
def add_signal_picture_content_slide(ppt:pptx.Presentation,data:dict,title:str,img,img_Dict:dict,key_message:str,title_color:RGBColor,big_title_color:RGBColor):
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    page = len(data)
    #添加标题
    title_box = slide.shapes.add_textbox(left=ppt.slide_width*3/5,top=ppt.slide_height/20,width=ppt.slide_width*2/5,height=ppt.slide_height/10)
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_box.text_frame.word_wrap = True
    title_box.text_frame.paragraphs[0].text = title
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.color.rgb = big_title_color
    #渐变色背景
    title_box.fill.gradient()
    title_box.fill.gradient_angle = 45
    title_box.fill.gradient_stops[0].color.rgb = title_color
    title_box.fill.gradient_stops[1].color.rgb = RGBColor(255,255,255)
    #添加图片
    if img != None:
        picture = slide.shapes.add_picture(img,left=ppt.slide_width*3/5,top=ppt.slide_height/5,width=ppt.slide_width*2/5,height=ppt.slide_height*8/10)
    #添加keyword
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,left=ppt.slide_width/40,top=ppt.slide_height*4/5,width=ppt.slide_width*3/5-ppt.slide_width/20,height=ppt.slide_height*4/25)
    shape.text = key_message
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    shape.text_frame.word_wrap = True
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.fill.solid()
    shape.fill.fore_color.rgb = title_color
    #添加主体内容
    left = ppt.slide_width/20
    height = ppt.slide_height/20
    width = ppt.slide_width*3/5-ppt.slide_width/10
    total_height = ppt.slide_height*3/4/page
    count=0
    for k,v in data.items():
        #标题
        little_title_box = slide.shapes.add_textbox(left=left,top=height+count*total_height,width=width,height=total_height*1/5)
        little_title_box.text_frame.paragraphs[0].text = k
        little_title_box.text_frame.paragraphs[0].font.bold = True
        little_title_box.text_frame.paragraphs[0].font.color.rgb = title_color
        little_title_box.text_frame.paragraphs[0].font.size = Pt(14)
        #内容
        content_box = slide.shapes.add_textbox(left=left,top=height+total_height*1/5+count*total_height,width=width,height=total_height*3/5)
        # content_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        content_box.text_frame.word_wrap = True
        content_box.text_frame.paragraphs[0].text = v
        content_box.text_frame.paragraphs[0].font.size = Pt(10)
        content_box.text_frame.paragraphs[0].line_spacing = 1.5
        count+=1
    #添加背景
    for k,v in img_Dict.items():
        picture = slide.shapes.add_picture(k,v[0],v[1],v[2],v[3])
        slide.shapes._spTree.remove(picture._element)
        slide.shapes._spTree.insert(2, picture._element)


def generate_powerpoint_presentation(
        parsed_data: dict,
        slides_template: str,
        template_color:str,
):
    """
    Create and save a PowerPoint presentation file containing the content in JSON format.
    parsed_data: 内容大纲
    slides_template: The path to the template PowerPoint presentation file.
    """
    #初始化页面

    mode_select = config.MODE_SELECT

    presentation = pptx.Presentation()
    ppt_tem = pptx.Presentation(slides_template)
    slide_width_inch, slide_height_inch = _get_slide_width_height_inches(ppt_tem)
    presentation.slide_width = Inches(slide_width_inch)
    presentation.slide_height = Inches(slide_height_inch)
    
    #按照模板制作PPT第一页
    full_imgDict = template.get_template_background(ppt_tem.slides[4],"full_img.jpg")
    slide_title = template.get_title_ppt(ppt_tem,parsed_data['title'],template.big_title_color[template_color])
    template.copy_template(presentation,slide_title)
    #按照模板填充目录
    slide_content = template.get_content_ppt(ppt_tem,parsed_data,slide_width_inch,slide_height_inch,template.small_content_color[template_color])
    template.copy_template(presentation,slide_content)
    #填充主体内容

    #并行制作PPT主体内容
    q = generate_ppt_stream.task_q
    excutor = ThreadPoolExecutor(max_workers=2)
    ppt_stream = io.BytesIO()
    def making_body():
        while True:
            # generate_ppt_stream.get_queue() 
            if not q.empty():
                #print("...开始获取.....")
                body = q.get()
                print(body)
                if isinstance(body,str):
                    if body == "<stop>":
                        #print("...结束循环.....")
                        slide_thank_you = template.get_thank_ppt(ppt_tem,template.big_title_color[template_color])
                        template.copy_template(presentation,slide_thank_you)    
                        presentation.save(ppt_stream)
                        ppt_stream.seek(0)
                        break
                i = body[1]
                j = body[2]               
                if j == 0:
                    #根据模板添加每一张开头页面
                    
                    slide_content_per = template.get_content_perppt(ppt_tem,parsed_data["chapter"][i]["heading"],i,template.big_content_color[template_color],template.big_title_color[template_color])
                    template.copy_template(presentation,slide_content_per)
                #添加每一小节内容
                point = body[0]
                title = point["heading"]
                items = point["bullet_points"]
                if len(items)<=2:
                    mode = random.randint(1,3)
                else:
                    mode = random.randint(1,4)
                
                #mode = 1使用左图右内容排版
                if mode == 1:
                    
                    if point["img_keywords"] != "":
                            points = point["img_keywords"].split(",")
                            rand = random.randint(0,len(points)-1)
                            if mode_select == 1:
                                #通过网络检索图片
                                img_key = ims.get_picture_from_url(query=points[rand])
                            elif mode_select == 2:
                                #通过HuggingFace API生成图片
                                img_keywords_en=generate_ppt_stream.translate_zn_to_en(points[rand])    
                                img_key = ims.get_picture_from_model_api_hug(query=img_keywords_en)
                            elif mode_select == 3:
                                #通过Zhipu API生成图片
                                img_key = ims.get_picture_from_model_api_zhipu(query=points[rand])
                            elif mode_select == 4:
                                img_keywords_en=generate_ppt_stream.translate_zn_to_en(points[rand])
                                img_key = ims.get_picture_from_model_huggingface(query=img_keywords_en)
                            if img_key == None:
                                for item in points:
                                    if mode_select == 1:
                                        #通过网络检索图片
                                        img_key = ims.get_picture_from_url(query=item)
                                    elif mode_select == 2:                            
                                        #通过HuggingFace API生成图片
                                        item_en = generate_ppt_stream.translate_zn_to_en(item)
                                        img_key = ims.get_picture_from_model_api_hug(query=item_en)
                                    elif mode_select == 3:
                                        #通过Zhipu API生成图片
                                        img_key = ims.get_picture_from_model_api_zhipu(query=item)
                                    elif mode_select == 4:
                                        item_en = generate_ppt_stream.translate_zn_to_en(item)
                                        img_key = ims.get_picture_from_model_huggingface(query=item_en)
                                    if img_key != None:
                                        break
                    add_signal_content_slide(presentation,items,title,img_key,full_imgDict,point["key_message"],template.small_content_color[template_color])
                elif mode == 2:
                    
                    mode_4 = random.randint(0,1)
                    # mode_4 = 0
                    if mode_4 == 0:
                        add_four_content_slide_no_picture(presentation,items,title,full_imgDict,point["key_message"],template.small_content_color[template_color])
                    else:
                        if len(items)<=4:
                            if point["img_keywords"] != "":
                                    points = point["img_keywords"].split(",")
                                    rand = random.randint(0,len(points)-1)
                                    if mode_select == 1:
                                        #通过网络检索图片
                                        img_key = ims.get_picture_from_url(query=points[rand])
                                    elif mode_select == 2:
                                        #通过HuggingFace API生成图片
                                        img_keywords_en=generate_ppt_stream.translate_zn_to_en(points[rand])    
                                        img_key = ims.get_picture_from_model_api_hug(query=img_keywords_en)
                                    elif mode_select == 3:
                                        #通过Zhipu API生成图片
                                        img_key = ims.get_picture_from_model_api_zhipu(query=points[rand])
                                    elif mode_select == 4:
                                        img_keywords_en=generate_ppt_stream.translate_zn_to_en(points[rand]) 
                                        img_key = ims.get_picture_from_model_huggingface(query=img_keywords_en)
                                    if img_key == None:
                                        for item in points:
                                            if mode_select == 1:
                                                #通过网络检索图片
                                                img_key = ims.get_picture_from_url(query=item)
                                            elif mode_select == 2:
                                                #通过HuggingFace API生成图片
                                                item_en = generate_ppt_stream.translate_zn_to_en(item)
                                                img_key = ims.get_picture_from_model_api_hug(query=item_en)
                                            elif mode_select == 3:
                                                img_key = ims.get_picture_from_model_api_zhipu(query=item)
                                            elif mode_select == 4:
                                                item_en = generate_ppt_stream.translate_zn_to_en(item)
                                                img_key = ims.get_picture_from_model_huggingface(query=item_en)
                                            if img_key != None:
                                                break
                            add_four_content_slide_with_signal_picture(presentation,items,img_key,icon_path,title,full_imgDict,point["key_message"],template.small_content_color[template_color])
                        else:
                            
                            imgs = []
                            for k,v in items.items():
                                if mode_select == 1:
                                    #通过网络检索图片
                                    img_key = ims.get_picture_from_url(query=k)
                                elif mode_select == 2:
                                    #通过HUGGINGFACE API 生成图片
                                    k_en = generate_ppt_stream.translate_zn_to_en(k)
                                    img_key = ims.get_picture_from_model_api_hug(query=k_en)
                                elif mode_select == 3:
                                    #通过zhipu API生成图片
                                    img_key = ims.get_picture_from_model_api_zhipu(k)
                                elif mode_select == 4:
                                    k_en = generate_ppt_stream.translate_zn_to_en(k)
                                    img_key = ims.get_picture_from_model_huggingface(k_en)
                                imgs.append(img_key)
                            add_content_slide_with_multiple_picture(presentation,items,imgs,title,full_imgDict,point["key_message"],template.small_content_color[template_color])
                elif mode == 3:
                    
                    if point["img_keywords"] != "":
                            points = point["img_keywords"].split(",")
                            rand = random.randint(0,len(points)-1)
                            if mode_select == 1:
                                #通过网络检索图片
                                img_key = ims.get_picture_from_url(query=points[rand])
                            elif mode_select == 2:
                                #通过HuggingFace API生成图片
                                img_keywords_en=generate_ppt_stream.translate_zn_to_en(points[rand])    
                                img_key = ims.get_picture_from_model_api_hug(query=img_keywords_en)
                            elif mode_select == 3:
                                #通过Zhipu API生成图片
                                img_key = ims.get_picture_from_model_api_zhipu(query=points[rand])
                            elif mode_select == 4:
                                img_keywords_en=generate_ppt_stream.translate_zn_to_en(points[rand])    
                                img_key = ims.get_picture_from_model_huggingface(query=img_keywords_en)
                            if img_key == None:
                                for item in points:
                                    if mode_select == 1:
                                        #通过网络检索图片
                                        img_key = ims.get_picture_from_url(query=item)
                                    elif mode_select == 2:
                                        #通过HuggingFace API生成图片
                                        item_en = generate_ppt_stream.translate_zn_to_en(item)
                                        img_key = ims.get_picture_from_model_api_hug(query=item_en)
                                    elif mode_select == 3:
                                        img_key = ims.get_picture_from_model_api_zhipu(query=item)
                                    elif mode_select == 4:
                                        item_en = generate_ppt_stream.translate_zn_to_en(item)
                                        img_key = ims.get_picture_from_model_huggingface(query=item_en)
                                    if img_key != None:
                                        break
                    add_signal_picture_content_slide(presentation,items,title,img_key,full_imgDict,point["key_message"],template.small_content_color[template_color],template.big_title_color[template_color])
                else:
                    
                    imgs = []
                    for k,v in items.items():
                        if mode_select == 1:
                            #通过网络检索图片
                            img_key = ims.get_picture_from_url(query=k)
                        elif mode_select == 2:
                            #通过HUGGINGFACE API 生成图片
                            k_en = generate_ppt_stream.translate_zn_to_en(k)
                            img_key = ims.get_picture_from_model_api_hug(query=k_en)
                        elif mode_select == 3:
                            #通过zhipu API生成图片
                            img_key = ims.get_picture_from_model_api_zhipu(k)
                        elif mode_select == 4:
                            k_en = generate_ppt_stream.translate_zn_to_en(k)
                            img_key = ims.get_picture_from_model_huggingface(k_en)
                        imgs.append(img_key)
                    add_content_slide_with_multiple_picture(presentation,items,imgs,title,full_imgDict,point["key_message"],template.small_content_color[template_color])

    excutor.submit(generate_ppt_stream.generate_each_slide_by_outline,parsed_data)
    task = excutor.submit(making_body) 
    while not task.done():
        continue
    return ppt_stream

# if __name__ == "__main__":
#     data = {
#     "title": "AI大模型：搞懂GGUF文件存储格式",
#     "chapter": [
#         {
#             "heading": "背景介绍",
#             "slides": [
#                 {
#                     "heading": "大模型概述",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "数据量与计算需求",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "现有存储格式问题",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "GGUF文件的引入",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 }
#             ]
#         },
#         {
#             "heading": "GGUF技术原理",
#             "slides": [
#                 {
#                     "heading": "设计目标与理念",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "结构化存储方式",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "压缩算法应用",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "跨平台兼容性",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 }
#             ]
#         },
#         {
#             "heading": "GGUF文件特性",
#             "slides": [
#                 {
#                     "heading": "格式优势分析",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "性能提升表现",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "兼容性与互操作性",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "实例解析：GGUF应用",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 }
#             ]
#         },
#         {
#             "heading": "实现细节与技术挑战",
#             "slides": [
#                 {
#                     "heading": "编码标准制定",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "数据解码过程",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "内存优化策略",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "性能测试与调优",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 }
#             ]
#         },
#         {
#             "heading": "实际应用案例",
#             "slides": [
#                 {
#                     "heading": "真实场景模拟演示",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "用户反馈与改进",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "案例分析：成功应用",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 },
#                 {
#                     "heading": "挑战与解决方案",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 }
#             ]
#         },
#         {
#             "heading": "总结",
#             "slides": [
#                 {
#                     "heading": "GGUF技术优势概述",
#                     "bullet_points": {},
#                     "key_message": "",
#                     "img_keywords": ""
#                 }
#             ]
#         }
#     ]
# }
#     generate_powerpoint_presentation(data,current_work_path,"青春")
    # 添加结束页
    # slide_thank_you = template.get_thank_ppt(ppt_tem,template.big_title_color[template_color])
    # template.copy_template(presentation,slide_thank_you)

    # ppt_stream = io.BytesIO()
    # presentation.save(ppt_stream)
    # ppt_stream.seek(0)

    # return ppt_stream
    # presentation.save("test.pptx")
    # return presentation


