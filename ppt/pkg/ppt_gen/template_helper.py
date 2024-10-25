import pptx
import pptx.presentation
from pptx.util import Cm,Pt,Inches
from pptx.dml.color import RGBColor
import copy

#大标题文字颜色
big_title_color = { "商务红": RGBColor(255,255,255),
               "商务灰": RGBColor(0,0,0),
               "青春":RGBColor(0,0,0),
               "绚烂":RGBColor(0,0,0),
               "水冷":RGBColor(0,0,0)

}
#目录大序号颜色
big_content_color = { "商务红": RGBColor(115,0,0),
               "商务灰": RGBColor(125,151,180),
               "青春":RGBColor(39,69,255),
               "绚烂":RGBColor(192,215,254),
               "水冷":RGBColor(143,166,255)

}
#目录小序号颜色
small_content_color = { "商务红": RGBColor(187,0,0),
               "商务灰": RGBColor(71,132,203),
               "青春":RGBColor(39,69,255),
               "绚烂":RGBColor(27,158,219),
               "水冷":RGBColor(39,69,255)

}

def get_template_background(slide,name:str):
    shapes = slide.shapes
    imgDict = {}
    for shape in shapes:
        if shape.name == "Image 0":
            with open(name,'wb') as f:
                f.write(shape.image.blob)
            imgDict[name] = [shape.left,shape.top,shape.width,shape.height]

    return imgDict
def copy_template(ppt:pptx.Presentation,slide):
    slides = ppt.slide_layouts[6]
    new_slide = ppt.slides.add_slide(slides)
    shapes = slide.shapes
    imgDict = {}
    for shape in shapes:
        if shape.name == "Image 0":
            with open(shape.name+'.jpg','wb') as f:
                f.write(shape.image.blob)

            imgDict[shape.name+'.jpg'] = [shape.left,shape.top,shape.width,shape.height]
        else:
            el = shape.element
            newel = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel,'p:extLst')
    for k,v in imgDict.items():
        picture = new_slide.shapes.add_picture(k,v[0],v[1],v[2],v[3])
        new_slide.shapes._spTree.remove(picture._element)
        new_slide.shapes._spTree.insert(2, picture._element)

    return True
def get_title_ppt(ppt:pptx.Presentation,data,title_color:RGBColor):
    """
        利用模板制作PPT封面

        ppt: 选用已有的模板
        data: 相关的内容
        slide: 返回制作好的PPT页面
    """
    slide = ppt.slides[0]
    shapes = slide.shapes
    for shape in shapes:
        if shape.name == "Text 1":
            shape.text_frame.paragraphs[0].text = data
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.size = Pt(37)
            shape.text_frame.paragraphs[0].font.color.rgb = title_color
        elif shape.name == "Text 0":
            shape.text = ""
        elif shape.name == "Text 3":
            shape.text = "By Yuanchat"
            shape.text_frame.paragraphs[0].font.color.rgb = title_color
    return slide
def get_content_ppt(ppt:pptx.Presentation,data,totalwidth,totalheight,title_color:RGBColor):
    """
        利用模板制作目录界面

        ppt: 选用已有的模板
        data: 相关的内容
        slide: 返回制作好的PPT页面

    """

    slide = ppt.slides[1]
    shapes = slide.shapes
    title = []
    for item in data["chapter"]:
        title.append(item["heading"])

    count = len(title)

    width = totalwidth/2
    per_height = totalheight/(count+1)
    height = 0
    for i in range(count):
        height += per_height
        textbox_n = shapes.add_textbox(left = Inches(width),top = Inches(height),height = Cm(1),width = Cm(1))
        textbox_t = shapes.add_textbox(left = Inches(width)+Inches(0.5),top = Inches(height),height = Cm(1),width = Cm(1))
        textbox_n.text_frame.paragraphs[0].text = str(i+1)
        textbox_n.text_frame.paragraphs[0].font.bold = True
        textbox_n.text_frame.paragraphs[0].font.size = Pt(28)
        textbox_n.text_frame.paragraphs[0].font.color.rgb = title_color
        textbox_t.text_frame.paragraphs[0].text = title[i]
        textbox_t.text_frame.paragraphs[0].font.bold = True
        textbox_t.text_frame.paragraphs[0].font.size = Pt(18)


    return slide
def get_content_perppt(ppt:pptx.Presentation,data:str,num:int,title_color:RGBColor,content_color:RGBColor):
    """
        利用模板制作每一章目录界面

        ppt: 选用已有的模板
        data: 相关的内容
        slide: 返回制作好的PPT页面
    """
    slide = ppt.slides[2]
    # title = []
    # for item in data["chapter"]:
    #     title.append(item["heading"])

    shapes = slide.shapes
    for shape in shapes:
        if shape.name == "Text 0":
            shape.text_frame.paragraphs[0].text = data
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.size = Pt(37)
            shape.text_frame.paragraphs[0].font.color.rgb = content_color
        elif shape.name == "Text 3":
            if num<9:
                shape.text_frame.paragraphs[0].text = "0"+str(num+1)
            else:
                shape.text_frame.paragraphs[0].text = str(num+1)
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.size = Pt(225)
            shape.text_frame.paragraphs[0].font.color.rgb = title_color
    return slide   
def get_thank_ppt(ppt:pptx.Presentation,title_color:RGBColor):
    """
        利用模板制作Thank页面

        ppt: 选用已有的模板
        data: 相关的内容
        slide: 返回制作好的PPT页面
    """
    slide = ppt.slides[5]
    shapes = slide.shapes
    for shape in shapes:
        if shape.name == "Text 0":
            shape.text_frame.paragraphs[0].text = "Thanks"
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.size = Pt(37)
            shape.text_frame.paragraphs[0].font.color.rgb = title_color
        elif shape.name == "Text 2":
            shape.text_frame.paragraphs[0].text = "Yuanchat"
            shape.text_frame.paragraphs[0].font.color.rgb = title_color

    return slide











